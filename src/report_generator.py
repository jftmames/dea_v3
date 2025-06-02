# src/report_generator.py

import pandas as pd
import datetime
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

def generate_html_report(
    df_dea: pd.DataFrame,
    df_tree: pd.DataFrame,
    df_eee: pd.DataFrame,
) -> str:
    """
    Genera un string con contenido HTML que incluye:
      - Título con fecha
      - Tabla de resultados DEA
      - Tabla de árbol de indagación (padre/hijo)
      - Tabla de metadatos EEE
    Retorna HTML listo para escribir a disco o servir como descarga.
    """
    fecha = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    html = (
        "<html><head><meta charset='utf-8'>"
        "<title>Reporte DEA Deliberativo</title></head><body>"
    )
    html += f"<h1>Reporte DEA Deliberativo – {fecha}</h1>"

    # Sección DEA
    html += "<h2>1. Resultados DEA</h2>"
    html += df_dea.to_html(index=False, border=1)

    # Sección Árbol
    html += "<h2>2. Estructura de Complejo de Indagación</h2>"
    html += df_tree.to_html(index=False, border=1)

    # Sección EEE
    html += "<h2>3. Métrico EEE y Metadatos</h2>"
    html += df_eee.to_html(index=False, border=1)

    html += "</body></html>"
    return html


def generate_excel_report(df_dea: pd.DataFrame, df_tree: pd.DataFrame, df_eee: pd.DataFrame) -> BytesIO:
    """
    Genera un archivo Excel en memoria con pestañas:
      - 'DEA Results'
      - 'Inquiry Tree'
      - 'EEE Metrics'
    Retorna un BytesIO conteniendo el Excel.
    Aplica formato de cabecera (bold, fondo, borde, centrado) y autoajusta ancho de columnas.
    """
    buffer = BytesIO()

    with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
        # Escribir cada DataFrame en su hoja
        df_dea.to_excel(writer, sheet_name="DEA Results", index=False)
        df_tree.to_excel(writer, sheet_name="Inquiry Tree", index=False)
        df_eee.to_excel(writer, sheet_name="EEE Metrics", index=False)

        workbook = writer.book
        header_fmt = workbook.add_format({
            "bold": True,
            "bg_color": "#D7E4BC",  # Un tono de verde claro para el fondo
            "border": 1,
            "align": "center",
            "valign": "vcenter" # Alineación vertical centrada
        })
        
        # Diccionario para mapear nombres de hoja a DataFrames originales
        hoja_a_df = {
            "DEA Results": df_dea,
            "Inquiry Tree": df_tree,
            "EEE Metrics": df_eee
        }

        # Iterar sobre cada hoja y su DataFrame original para aplicar formato
        for sheet_name, df_original in hoja_a_df.items():
            worksheet = writer.sheets[sheet_name]
            
            # Aplicar formato a las cabeceras usando los nombres de columna del DataFrame original
            for col_num, col_name in enumerate(df_original.columns):
                worksheet.write(0, col_num, col_name, header_fmt)
            
            # Autoajustar el ancho de las columnas
            for col_num, col_name in enumerate(df_original.columns):
                max_len = max(df_original[col_name].astype(str).map(len).max(), len(col_name)) + 2 # +2 para un pequeño margen
                worksheet.set_column(col_num, col_num, max_len)

    buffer.seek(0)
    return buffer


def generate_pptx_report(
    df_dea: pd.DataFrame,
    df_tree: pd.DataFrame,
    df_eee: pd.DataFrame
) -> BytesIO:
    """
    Genera un archivo PPTX en memoria con slides para:
      1. Resultados DEA (tabla)
      2. Árbol de indagación (tabla)
      3. Métricas EEE (tabla)
    Retorna un BytesIO conteniendo el PPTX.
    """
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # layout en blanco

    def add_table_slide(df: pd.DataFrame, title: str):
        slide = prs.slides.add_slide(blank_slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title # Asignar el título

        rows, cols = df.shape
        # Ajustar el tamaño y la posición de la tabla para que quepa en la diapositiva
        # Si la tabla es muy grande, podría necesitar un ajuste más dinámico o paginación
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(min(Inches(7), Inches(0.4 * (rows + 1)))) # Ajustar altura máxima para evitar que se salga

        table_shape = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
        table = table_shape.table

        # Estilo de la tabla (opcional, se puede aplicar un estilo predefinido)
        # table.cell(0, 0)._tc.xmlL_cell(0, 0) for cell in table.iter_cells()
        # for cell in table.iter_cells():
        #     cell.fill.solid()
        #     cell.fill.fore_color.rgb = RGBColor(0xDAE3F3) # Ejemplo de color de fondo claro
        #     cell.text_frame.paragraphs[0].font.size = Pt(10)

        # Escribir cabeceras
        for j, col_name in enumerate(df.columns):
            cell = table.cell(0, j)
            cell.text = str(col_name)
            # Opcional: Estilo para la cabecera
            # cell.fill.solid()
            # cell.fill.fore_color.rgb = RGBColor(0xBFBFBF) # Gris oscuro
            # cell.text_frame.paragraphs[0].font.bold = True
            # cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


        # Escribir datos
        for i in range(rows):
            for j, col_name in enumerate(df.columns):
                cell = table.cell(i + 1, j)
                cell.text = str(df.iloc[i, j])
                # cell.text_frame.paragraphs[0].font.size = Pt(9)

        # Autoajustar ancho de columnas en PowerPoint es más complejo y a menudo
        # requiere iteración o cálculo manual basado en contenido.
        # Por ahora, se dejarán los anchos por defecto o se ajustará manualmente en PPT.
        # table.columns[j].width = Inches(...) # Esto sería si quisiéramos fijar un ancho

    # Añadir slides para cada sección
    add_table_slide(df_dea, "1. Resultados DEA")
    add_table_slide(df_tree, "2. Árbol de Indagación")
    add_table_slide(df_eee, "3. Métricas EEE")

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
